#!/usr/bin/env python3
# This code was developed with the help of AI.

from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import sys
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Set

import fitz
import pytesseract
import xlrd
import xlwt
from docx import Document
from docx.shared import Inches as DocxInches
from openpyxl import Workbook, load_workbook
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches as PptxInches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


class Logger:
    @staticmethod
    def info(message: str) -> None:
        print(f"INFO: {message}", file=sys.stderr)

    @staticmethod
    def warning(message: str) -> None:
        print(f"WARNING: {message}", file=sys.stderr)

    @staticmethod
    def error(message: str) -> None:
        print(f"ERROR: {message}", file=sys.stderr)


class PythonVersionEnforcer:
    def __init__(self, required_version: str):
        self.required_version = required_version

    def validate(self) -> None:
        parts = self.required_version.strip().split(".")
        if len(parts) != 2 or not all(part.isdigit() for part in parts):
            raise RuntimeError(
                f"Invalid required Python version format: '{self.required_version}'. Use '<major>.<minor>' like '3.12'."
            )

        required_major, required_minor = int(parts[0]), int(parts[1])
        current = (sys.version_info.major, sys.version_info.minor)
        if current != (required_major, required_minor):
            raise RuntimeError(
                f"Python {required_major}.{required_minor} is required, but current runtime is {current[0]}.{current[1]}."
            )


@dataclass(frozen=True)
class AppConfig:
    input_dir: Path
    output_dir: Path
    output_jsonl: Path
    libreoffice_bin: Optional[str]
    tesseract_bin: str
    supported_extensions: Set[str]
    required_python: str
    sample_image_name: str
    sample_pdf_name: str
    sample_docx_name: str
    sample_xlsx_name: str
    sample_xls_name: str
    sample_pptx_name: str
    sample_ppt_name: str


class AppConfigLoader:
    ENV_INPUT_DIR = "DOC_PREPROCESS_EXAMPLES_DIR"
    ENV_OUTPUT_DIR = "DOC_PREPROCESS_OUTPUT_DIR"
    ENV_OUTPUT_JSONL = "DOC_PREPROCESS_OUTPUT_JSONL"
    ENV_LIBREOFFICE_BIN = "DOC_PREPROCESS_LIBREOFFICE_BIN"
    ENV_TESSERACT_BIN = "DOC_PREPROCESS_TESSERACT_BIN"
    ENV_SUPPORTED_EXTENSIONS = "DOC_PREPROCESS_SUPPORTED_EXTENSIONS"
    ENV_REQUIRED_PYTHON = "DOC_PREPROCESS_REQUIRED_PYTHON"

    ENV_SAMPLE_IMAGE_NAME = "DOC_PREPROCESS_SAMPLE_IMAGE_NAME"
    ENV_SAMPLE_PDF_NAME = "DOC_PREPROCESS_SAMPLE_PDF_NAME"
    ENV_SAMPLE_DOCX_NAME = "DOC_PREPROCESS_SAMPLE_DOCX_NAME"
    ENV_SAMPLE_XLSX_NAME = "DOC_PREPROCESS_SAMPLE_XLSX_NAME"
    ENV_SAMPLE_XLS_NAME = "DOC_PREPROCESS_SAMPLE_XLS_NAME"
    ENV_SAMPLE_PPTX_NAME = "DOC_PREPROCESS_SAMPLE_PPTX_NAME"
    ENV_SAMPLE_PPT_NAME = "DOC_PREPROCESS_SAMPLE_PPT_NAME"

    def __init__(self, env: Dict[str, str], args: argparse.Namespace):
        self.env = env
        self.args = args

    def load(self) -> AppConfig:
        input_dir = self._path_value(self.args.input_dir, self.ENV_INPUT_DIR)
        output_dir = self._path_value(self.args.output_dir, self.ENV_OUTPUT_DIR)
        output_jsonl = self._path_value(self.args.output_jsonl, self.ENV_OUTPUT_JSONL)

        supported_extensions = self._extensions_value(
            self._text_value(self.args.supported_extensions, self.ENV_SUPPORTED_EXTENSIONS)
        )
        required_python = self._text_value(self.args.required_python, self.ENV_REQUIRED_PYTHON)

        libreoffice_bin_raw = self._optional_text_value(self.args.libreoffice_bin, self.ENV_LIBREOFFICE_BIN)
        libreoffice_bin = libreoffice_bin_raw if libreoffice_bin_raw else None

        tesseract_bin = self._text_value(self.args.tesseract_bin, self.ENV_TESSERACT_BIN)

        return AppConfig(
            input_dir=input_dir,
            output_dir=output_dir,
            output_jsonl=output_jsonl,
            libreoffice_bin=libreoffice_bin,
            tesseract_bin=tesseract_bin,
            supported_extensions=supported_extensions,
            required_python=required_python,
            sample_image_name=self._text_value(self.args.sample_image_name, self.ENV_SAMPLE_IMAGE_NAME),
            sample_pdf_name=self._text_value(self.args.sample_pdf_name, self.ENV_SAMPLE_PDF_NAME),
            sample_docx_name=self._text_value(self.args.sample_docx_name, self.ENV_SAMPLE_DOCX_NAME),
            sample_xlsx_name=self._text_value(self.args.sample_xlsx_name, self.ENV_SAMPLE_XLSX_NAME),
            sample_xls_name=self._text_value(self.args.sample_xls_name, self.ENV_SAMPLE_XLS_NAME),
            sample_pptx_name=self._text_value(self.args.sample_pptx_name, self.ENV_SAMPLE_PPTX_NAME),
            sample_ppt_name=self._text_value(self.args.sample_ppt_name, self.ENV_SAMPLE_PPT_NAME),
        )

    def _path_value(self, arg_value: Optional[str], env_name: str) -> Path:
        raw = self._text_value(arg_value, env_name)
        return Path(raw).resolve()

    def _optional_text_value(self, arg_value: Optional[str], env_name: str) -> str:
        if arg_value is not None and str(arg_value).strip() != "":
            return str(arg_value).strip()
        return str(self.env.get(env_name, "")).strip()

    def _text_value(self, arg_value: Optional[str], env_name: str) -> str:
        value = self._optional_text_value(arg_value, env_name)
        if value == "":
            raise RuntimeError(f"Missing configuration. Set environment variable '{env_name}' or pass CLI option.")
        return value

    @staticmethod
    def _extensions_value(raw_extensions: str) -> Set[str]:
        parsed = []
        for entry in raw_extensions.split(","):
            item = entry.strip().lower()
            if not item:
                continue
            parsed.append(item if item.startswith(".") else f".{item}")

        if not parsed:
            raise RuntimeError("No supported extensions configured.")
        return set(parsed)


class CLIParserFactory:
    @staticmethod
    def build() -> argparse.ArgumentParser:
        parser = argparse.ArgumentParser(description="Generate examples and preprocess documents to JSONL.")
        parser.add_argument("--input-dir", default=None, type=str)
        parser.add_argument("--output-dir", default=None, type=str)
        parser.add_argument("--output-jsonl", default=None, type=str)
        parser.add_argument("--libreoffice-bin", default=None, type=str)
        parser.add_argument("--tesseract-bin", default=None, type=str)
        parser.add_argument("--supported-extensions", default=None, type=str)
        parser.add_argument("--required-python", default=None, type=str)

        parser.add_argument("--sample-image-name", default=None, type=str)
        parser.add_argument("--sample-pdf-name", default=None, type=str)
        parser.add_argument("--sample-docx-name", default=None, type=str)
        parser.add_argument("--sample-xlsx-name", default=None, type=str)
        parser.add_argument("--sample-xls-name", default=None, type=str)
        parser.add_argument("--sample-pptx-name", default=None, type=str)
        parser.add_argument("--sample-ppt-name", default=None, type=str)
        return parser


class PathManager:
    @staticmethod
    def ensure_dir(path: Path) -> None:
        path.mkdir(parents=True, exist_ok=True)


class LibreOfficeService:
    def __init__(self, libreoffice_bin: Optional[str]):
        self.libreoffice_bin = libreoffice_bin

    def is_available(self) -> bool:
        return bool(self.libreoffice_bin)

    def convert(self, source_file: Path, target_format: str, out_dir: Path) -> Path:
        if not self.libreoffice_bin:
            raise RuntimeError("LibreOffice binary is not configured.")

        PathManager.ensure_dir(out_dir)
        target_file = out_dir / f"{source_file.stem}.{target_format.lower()}"
        if target_file.exists():
            target_file.unlink()

        command = [
            self.libreoffice_bin,
            "--headless",
            "--convert-to",
            target_format,
            "--outdir",
            str(out_dir),
            str(source_file),
        ]
        process = subprocess.run(command, capture_output=True, text=True)
        if process.returncode != 0:
            stderr = (process.stderr or "").strip()
            stdout = (process.stdout or "").strip()
            raise RuntimeError(
                f"LibreOffice conversion failed for {source_file.name}. stdout={stdout} stderr={stderr}"
            )

        if target_file.exists():
            return target_file

        for candidate in out_dir.glob(f"{source_file.stem}.*"):
            if candidate.suffix.lower() == f".{target_format.lower()}":
                return candidate

        raise RuntimeError(f"LibreOffice reported success but no converted file was found for {source_file.name}")


class OCRService:
    def __init__(self, tesseract_bin: str):
        self.tesseract_bin = tesseract_bin
        self._available = bool(shutil.which(self.tesseract_bin))
        if self._available:
            pytesseract.pytesseract.tesseract_cmd = self.tesseract_bin

    @property
    def available(self) -> bool:
        return self._available

    def extract_from_image_bytes(self, image_bytes: bytes) -> Optional[str]:
        if not self._available:
            return None
        try:
            with Image.open(BytesIO(image_bytes)) as img:
                text = pytesseract.image_to_string(img).strip()
            return text or None
        except Exception:
            return None


class SampleFileGenerator:
    def __init__(self, config: AppConfig, libreoffice_service: LibreOfficeService):
        self.config = config
        self.libreoffice_service = libreoffice_service

    def generate(self) -> None:
        PathManager.ensure_dir(self.config.input_dir)

        image_path = self.config.input_dir / self.config.sample_image_name
        pdf_path = self.config.input_dir / self.config.sample_pdf_name
        docx_path = self.config.input_dir / self.config.sample_docx_name
        xlsx_path = self.config.input_dir / self.config.sample_xlsx_name
        xls_path = self.config.input_dir / self.config.sample_xls_name
        pptx_path = self.config.input_dir / self.config.sample_pptx_name
        ppt_path = self.config.input_dir / self.config.sample_ppt_name

        self._create_ocr_image(image_path)
        self._create_sample_pdf(pdf_path, image_path)
        self._create_sample_docx(docx_path, image_path)
        self._create_sample_xlsx(xlsx_path)
        self._create_sample_xls(xls_path)
        self._create_sample_pptx(pptx_path, image_path)

        if self.libreoffice_service.is_available():
            try:
                self.libreoffice_service.convert(pptx_path, "ppt", self.config.input_dir)
            except Exception as exc:
                Logger.warning(f"Could not generate legacy .ppt example: {exc}")
        else:
            if ppt_path.exists():
                ppt_path.unlink()
            Logger.info("Legacy .ppt example was not generated because LibreOffice is not configured.")

    @staticmethod
    def _create_ocr_image(image_path: Path) -> None:
        img = Image.new("RGB", (900, 220), "white")
        draw = ImageDraw.Draw(img)
        draw.text((30, 80), "OCR DEMO TEXT 12345", fill="black")
        draw.text((30, 130), "Doc preprocessing sample image.", fill="black")
        img.save(image_path)

    @staticmethod
    def _create_sample_pdf(pdf_path: Path, image_path: Path) -> None:
        pdf_canvas = canvas.Canvas(str(pdf_path), pagesize=letter)
        pdf_canvas.setFont("Helvetica", 12)
        pdf_canvas.drawString(72, 750, "Sample PDF - Page 1")
        pdf_canvas.drawString(72, 730, "This page includes native text and an embedded image.")
        pdf_canvas.drawImage(str(image_path), 72, 520, width=420, height=120, preserveAspectRatio=True, mask="auto")
        pdf_canvas.showPage()
        pdf_canvas.setFont("Helvetica", 12)
        pdf_canvas.drawString(72, 750, "Sample PDF - Page 2")
        pdf_canvas.drawString(72, 730, "Second page for unit-wise page extraction.")
        pdf_canvas.save()

    @staticmethod
    def _create_sample_docx(docx_path: Path, image_path: Path) -> None:
        document = Document()
        document.add_heading("Sample DOCX", level=1)
        document.add_paragraph("This DOCX demonstrates document-level extraction.")
        document.add_paragraph("It includes paragraphs and one image for OCR best-effort.")
        document.add_picture(str(image_path), width=DocxInches(5.5))
        document.save(str(docx_path))

    @staticmethod
    def _create_sample_xlsx(xlsx_path: Path) -> None:
        workbook = Workbook()
        overview = workbook.active
        overview.title = "Overview"
        overview.append(["Key", "Value"])
        overview.append(["Project", "Doc preprocessing"])
        overview.append(["Status", "Sample data"])

        metrics = workbook.create_sheet("Metrics")
        metrics.append(["Month", "Count"])
        metrics.append(["January", 10])
        metrics.append(["February", 17])
        workbook.save(str(xlsx_path))

    @staticmethod
    def _create_sample_xls(xls_path: Path) -> None:
        workbook = xlwt.Workbook()

        sheet1 = workbook.add_sheet("Overview")
        rows_1 = [
            ["Key", "Value"],
            ["Legacy", "XLS format"],
            ["Status", "Sample data"],
        ]
        for row_idx, row in enumerate(rows_1):
            for col_idx, value in enumerate(row):
                sheet1.write(row_idx, col_idx, value)

        sheet2 = workbook.add_sheet("Metrics")
        rows_2 = [
            ["Quarter", "Revenue"],
            ["Q1", 100],
            ["Q2", 120],
        ]
        for row_idx, row in enumerate(rows_2):
            for col_idx, value in enumerate(row):
                sheet2.write(row_idx, col_idx, value)

        workbook.save(str(xls_path))

    @staticmethod
    def _create_sample_pptx(pptx_path: Path, image_path: Path) -> None:
        presentation = Presentation()

        slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide1.shapes.title.text = "Sample PPTX - Slide 1"
        slide1.placeholders[1].text = "This slide contains text for slide-wise extraction."

        slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])
        title_shape = slide2.shapes.title
        if title_shape is not None:
            title_shape.text = "Sample PPTX - Slide 2"

        text_box = slide2.shapes.add_textbox(PptxInches(0.8), PptxInches(1.4), PptxInches(8.0), PptxInches(1.2))
        text_box.text_frame.text = "This slide includes an image for OCR best-effort."
        slide2.shapes.add_picture(str(image_path), PptxInches(0.8), PptxInches(2.2), width=PptxInches(7.5))

        presentation.save(str(pptx_path))


class MarkdownFormatter:
    @staticmethod
    def to_text(value: Any) -> str:
        if value is None:
            return ""
        if isinstance(value, float) and value.is_integer():
            return str(int(value))
        return str(value)

    def sheet_rows_to_markdown(self, title: str, rows: List[List[Any]]) -> str:
        cleaned: List[List[str]] = []
        for row in rows:
            converted = [self.to_text(v).strip() for v in row]
            if any(cell != "" for cell in converted):
                cleaned.append(converted)

        if not cleaned:
            return f"# Sheet {title}\n\n_No text extracted._"

        max_cols = max(len(row) for row in cleaned)
        headers = [f"C{i}" for i in range(1, max_cols + 1)]

        lines = [
            f"# Sheet {title}",
            "",
            "|" + "|".join(headers) + "|",
            "|" + "|".join(["---"] * max_cols) + "|",
        ]

        for row in cleaned:
            padded = row + [""] * (max_cols - len(row))
            escaped = [cell.replace("|", "\\|").replace("\n", "<br>") for cell in padded]
            lines.append("|" + "|".join(escaped) + "|")

        return "\n".join(lines)


class RecordFactory:
    @staticmethod
    def build(
        source_file: Path,
        input_format: str,
        unit_number: int,
        unit_type: str,
        text_markdown: str,
        ocr_image_text: Optional[str],
        metadata: Dict[str, Any],
    ) -> Dict[str, Any]:
        return {
            "source_file_path": str(source_file.resolve()),
            "input_format": input_format,
            "unit_number": unit_number,
            "unit_type": unit_type,
            "text_markdown": text_markdown,
            "ocr_image_text": ocr_image_text,
            "metadata": metadata,
        }


class DocumentPreprocessor:
    def __init__(
        self,
        config: AppConfig,
        ocr_service: OCRService,
        libreoffice_service: LibreOfficeService,
        markdown_formatter: MarkdownFormatter,
    ):
        self.config = config
        self.ocr_service = ocr_service
        self.libreoffice_service = libreoffice_service
        self.markdown_formatter = markdown_formatter

    def preprocess_all(self) -> None:
        PathManager.ensure_dir(self.config.output_dir)

        docs = sorted(
            path
            for path in self.config.input_dir.rglob("*")
            if path.is_file() and path.suffix.lower() in self.config.supported_extensions
        )

        if not docs:
            raise RuntimeError(
                "No supported documents found in input directory. "
                f"Configured extensions: {sorted(self.config.supported_extensions)}"
            )

        errors: List[str] = []
        with self.config.output_jsonl.open("w", encoding="utf-8") as output_handle:
            for doc_file in docs:
                try:
                    records = self._process_file(doc_file)
                    for record in records:
                        output_handle.write(json.dumps(record, ensure_ascii=False) + "\n")
                except Exception as exc:
                    errors.append(f"{doc_file}: {exc}")

        if errors:
            for error in errors:
                Logger.error(error)
            raise RuntimeError("One or more files failed during preprocessing.")

    def _process_file(self, file_path: Path) -> List[Dict[str, Any]]:
        suffix = file_path.suffix.lower()
        if suffix == ".pdf":
            return self._process_pdf(file_path)
        if suffix == ".docx":
            return self._process_docx(file_path)
        if suffix == ".xlsx":
            return self._process_xlsx(file_path)
        if suffix == ".xls":
            return self._process_xls(file_path)
        if suffix in {".pptx", ".ppt"}:
            return self._process_ppt(file_path)
        raise RuntimeError(f"Unsupported file type: {file_path}")

    def _extract_pdf_docling_markdown(self, pdf_path: Path) -> Dict[str, Any]:
        try:
            from docling.document_converter import DocumentConverter, PdfFormatOption
            from docling.datamodel.base_models import InputFormat
        except Exception as exc:
            return {
                "ok": False,
                "pages": {},
                "error": f"Docling import failed: {exc}",
                "docling_ocr_enabled": False,
            }

        try:
            docling_ocr_enabled = False
            try:
                from docling.datamodel.pipeline_options import PdfPipelineOptions

                options = PdfPipelineOptions()
                if hasattr(options, "do_ocr"):
                    setattr(options, "do_ocr", True)
                    docling_ocr_enabled = True
                converter = DocumentConverter(
                    format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=options)}
                )
            except Exception:
                converter = DocumentConverter(format_options={InputFormat.PDF: PdfFormatOption()})

            result = converter.convert(str(pdf_path))
            document = result.document

            page_count = int(document.num_pages()) if hasattr(document, "num_pages") else 0
            if page_count < 1:
                with fitz.open(str(pdf_path)) as pdf_document:
                    page_count = pdf_document.page_count

            pages: Dict[int, str] = {}
            for page_number in range(1, page_count + 1):
                text_md = ""
                if hasattr(document, "export_to_markdown"):
                    try:
                        text_md = document.export_to_markdown(page_no=page_number) or ""
                    except TypeError:
                        if page_number == 1:
                            text_md = document.export_to_markdown() or ""
                pages[page_number] = text_md

            return {
                "ok": True,
                "pages": pages,
                "error": "",
                "docling_ocr_enabled": docling_ocr_enabled,
            }
        except Exception as exc:
            return {
                "ok": False,
                "pages": {},
                "error": f"Docling conversion failed: {exc}",
                "docling_ocr_enabled": False,
            }

    def _extract_pdf_ocr_by_page(self, pdf_path: Path) -> Dict[int, str]:
        if not self.ocr_service.available:
            return {}

        page_ocr: Dict[int, str] = {}
        with fitz.open(str(pdf_path)) as pdf_document:
            for idx, page in enumerate(pdf_document, start=1):
                try:
                    pixmap = page.get_pixmap(dpi=220)
                    text = self.ocr_service.extract_from_image_bytes(pixmap.tobytes("png"))
                    if text:
                        page_ocr[idx] = text
                except Exception:
                    continue
        return page_ocr

    def _process_pdf(self, file_path: Path) -> List[Dict[str, Any]]:
        records: List[Dict[str, Any]] = []
        docling_data = self._extract_pdf_docling_markdown(file_path)
        ocr_by_page = self._extract_pdf_ocr_by_page(file_path)

        with fitz.open(str(file_path)) as pdf_document:
            total_pages = pdf_document.page_count
            for page_no in range(1, total_pages + 1):
                text_markdown = ""
                if docling_data["ok"]:
                    text_markdown = (docling_data["pages"].get(page_no) or "").strip()

                if not text_markdown:
                    native_text = pdf_document[page_no - 1].get_text("text").strip()
                    text_markdown = (
                        f"# Page {page_no}\n\n{native_text}" if native_text else f"# Page {page_no}\n\n_No text extracted._"
                    )

                records.append(
                    RecordFactory.build(
                        source_file=file_path,
                        input_format="pdf",
                        unit_number=page_no,
                        unit_type="page",
                        text_markdown=text_markdown,
                        ocr_image_text=ocr_by_page.get(page_no),
                        metadata={
                            "total_pages": total_pages,
                            "page_number": page_no,
                            "docling_used": bool(docling_data["ok"]),
                            "docling_ocr_enabled": bool(docling_data.get("docling_ocr_enabled", False)),
                            "docling_error": docling_data.get("error", "") if not docling_data["ok"] else "",
                            "tesseract_available": self.ocr_service.available,
                        },
                    )
                )

        return records

    def _process_docx(self, file_path: Path) -> List[Dict[str, Any]]:
        document = Document(str(file_path))
        paragraphs = [p.text.strip() for p in document.paragraphs if p.text and p.text.strip()]
        markdown_text = "# Document\n\n" + ("\n\n".join(paragraphs) if paragraphs else "_No text extracted._")

        ocr_texts: List[str] = []
        image_count = 0
        for relation in document.part.rels.values():
            if "image" in relation.reltype:
                image_count += 1
                text = self.ocr_service.extract_from_image_bytes(relation.target_part.blob)
                if text:
                    ocr_texts.append(text)

        return [
            RecordFactory.build(
                source_file=file_path,
                input_format="docx",
                unit_number=1,
                unit_type="document",
                text_markdown=markdown_text,
                ocr_image_text="\n\n".join(ocr_texts) if ocr_texts else None,
                metadata={
                    "paragraph_count": len(paragraphs),
                    "image_count": image_count,
                    "tesseract_available": self.ocr_service.available,
                },
            )
        ]

    def _ensure_pptx_for_processing(self, file_path: Path) -> Path:
        if file_path.suffix.lower() == ".pptx":
            return file_path

        if not self.libreoffice_service.is_available():
            raise RuntimeError(
                "Cannot process .ppt without LibreOffice. Configure DOC_PREPROCESS_LIBREOFFICE_BIN or install soffice."
            )

        converted_dir = self.config.output_dir / "_converted"
        return self.libreoffice_service.convert(file_path, "pptx", converted_dir)

    def _process_ppt(self, file_path: Path) -> List[Dict[str, Any]]:
        effective_pptx = self._ensure_pptx_for_processing(file_path)
        presentation = Presentation(str(effective_pptx))

        records: List[Dict[str, Any]] = []
        total_slides = len(presentation.slides)

        for slide_index, slide in enumerate(presentation.slides, start=1):
            text_blocks: List[str] = []
            ocr_blocks: List[str] = []
            image_count = 0

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = (shape.text or "").strip()
                    if text:
                        text_blocks.append(text)

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                    try:
                        extracted = self.ocr_service.extract_from_image_bytes(shape.image.blob)
                        if extracted:
                            ocr_blocks.append(extracted)
                    except Exception:
                        pass

            body = "\n\n".join(text_blocks) if text_blocks else "_No text extracted._"
            markdown_text = f"# Slide {slide_index}\n\n{body}"

            records.append(
                RecordFactory.build(
                    source_file=file_path,
                    input_format=file_path.suffix.lower().lstrip("."),
                    unit_number=slide_index,
                    unit_type="slide",
                    text_markdown=markdown_text,
                    ocr_image_text="\n\n".join(ocr_blocks) if ocr_blocks else None,
                    metadata={
                        "total_slides": total_slides,
                        "slide_number": slide_index,
                        "image_count": image_count,
                        "processed_from": str(effective_pptx.resolve()),
                        "tesseract_available": self.ocr_service.available,
                    },
                )
            )

        return records

    def _process_xlsx(self, file_path: Path) -> List[Dict[str, Any]]:
        workbook = load_workbook(str(file_path), data_only=True)
        records: List[Dict[str, Any]] = []
        total_sheets = len(workbook.worksheets)

        for idx, worksheet in enumerate(workbook.worksheets, start=1):
            rows = [list(row) for row in worksheet.iter_rows(values_only=True)]
            text_markdown = self.markdown_formatter.sheet_rows_to_markdown(worksheet.title, rows)

            ocr_blocks: List[str] = []
            image_count = 0
            for image in getattr(worksheet, "_images", []):
                image_count += 1
                try:
                    extracted = self.ocr_service.extract_from_image_bytes(image._data())
                    if extracted:
                        ocr_blocks.append(extracted)
                except Exception:
                    pass

            records.append(
                RecordFactory.build(
                    source_file=file_path,
                    input_format="xlsx",
                    unit_number=idx,
                    unit_type="sheet",
                    text_markdown=text_markdown,
                    ocr_image_text="\n\n".join(ocr_blocks) if ocr_blocks else None,
                    metadata={
                        "sheet_name": worksheet.title,
                        "total_sheets": total_sheets,
                        "sheet_number": idx,
                        "image_count": image_count,
                        "tesseract_available": self.ocr_service.available,
                    },
                )
            )

        return records

    def _process_xls(self, file_path: Path) -> List[Dict[str, Any]]:
        workbook = xlrd.open_workbook(str(file_path))
        records: List[Dict[str, Any]] = []
        total_sheets = workbook.nsheets

        for idx in range(total_sheets):
            sheet = workbook.sheet_by_index(idx)
            rows: List[List[Any]] = []
            for row_index in range(sheet.nrows):
                rows.append([sheet.cell_value(row_index, col_idx) for col_idx in range(sheet.ncols)])

            text_markdown = self.markdown_formatter.sheet_rows_to_markdown(sheet.name, rows)
            records.append(
                RecordFactory.build(
                    source_file=file_path,
                    input_format="xls",
                    unit_number=idx + 1,
                    unit_type="sheet",
                    text_markdown=text_markdown,
                    ocr_image_text=None,
                    metadata={
                        "sheet_name": sheet.name,
                        "total_sheets": total_sheets,
                        "sheet_number": idx + 1,
                        "image_count": 0,
                        "tesseract_available": self.ocr_service.available,
                    },
                )
            )

        return records


class PreprocessApplication:
    def __init__(self, args: argparse.Namespace, env: Dict[str, str]):
        self.config = AppConfigLoader(env=env, args=args).load()

    def run(self) -> int:
        PythonVersionEnforcer(self.config.required_python).validate()

        PathManager.ensure_dir(self.config.input_dir)
        PathManager.ensure_dir(self.config.output_dir)

        libreoffice_service = LibreOfficeService(self.config.libreoffice_bin)
        ocr_service = OCRService(self.config.tesseract_bin)

        if not ocr_service.available:
            Logger.warning(
                f"OCR binary '{self.config.tesseract_bin}' was not found. OCR image text will be empty where OCR is required."
            )

        generator = SampleFileGenerator(config=self.config, libreoffice_service=libreoffice_service)
        generator.generate()

        processor = DocumentPreprocessor(
            config=self.config,
            ocr_service=ocr_service,
            libreoffice_service=libreoffice_service,
            markdown_formatter=MarkdownFormatter(),
        )
        processor.preprocess_all()

        Logger.info(f"Wrote JSONL output to {self.config.output_jsonl}")
        return 0


class ApplicationEntryPoint:
    @staticmethod
    def run() -> int:
        parser = CLIParserFactory.build()
        args = parser.parse_args()
        app = PreprocessApplication(args=args, env=dict(os.environ))
        return app.run()


if __name__ == "__main__":
    try:
        raise SystemExit(ApplicationEntryPoint.run())
    except Exception as exc:
        Logger.error(str(exc))
        raise SystemExit(1)
